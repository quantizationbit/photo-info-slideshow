#!/usr/bin/env python3
"""
process_photos.py — Incremental trip photo processing pipeline.

Trip name, output filenames, and the trip's date range (for "Day N" tags) are
all auto-determined from the folder the script sits in — copy this file into
any trip folder and run it, no per-trip editing required. See TRIP_NAME /
TRIP_SLUG below, and TRIP_START_DATE / TRIP_END_DATE which are computed from
the photos' own EXIF timestamps at runtime.

For each new/changed source photo:
  - Extracts EXIF (GPS, date/time, camera fingerprint)
  - Looks up USGS elevation
  - Reverse-geocodes nearest area (Nominatim zoom=17)
  - Generates _topo.jpg (USGS topo map, 1000x1000, 1000-yard radius)
  - Generates _fullinfo.jpg (caption burn-in + photographer credit;
    topo inset only included with --burn-in)
  - Generates _nearestarea.jpg (area phrase only)

Manifest (<slug>_Photo_Manifest.docx) is always rebuilt from the state cache,
so unchanged photos cost no API calls.

Slideshow (<slug>_Slideshow.mp4): 4K UHD H.264 CRF=14. Script always prompts
before encoding since it is time-consuming. Sequence per photo:
  topo (4s, black-bg, skipped if closer than an adaptive per-trip distance
    threshold — median inter-photo gap × TOPO_ADAPTIVE_MULT, or TOPO_SKIP_DIST
    as a fallback — from the previous photo's topo),
  fullinfo (3s, black-bg fill if not 16:9),
  nearestarea (1s, black-bg fill if not 16:9),
  original (2s, black-bg fill if not 16:9).

Usage:
  python3 process_photos.py              # incremental (skip unchanged), then prompt for slideshow
  python3 process_photos.py --force      # reprocess all photos
  python3 process_photos.py --manifest   # rebuild manifest only, no image processing
  python3 process_photos.py --slideshow  # build slideshow without prompting
  python3 process_photos.py --no-slideshow  # skip slideshow prompt entirely
  python3 process_photos.py --burn-in    # burn topo inset into _fullinfo images
  python3 process_photos.py --powerpoint # also build a black-background .pptx slide deck
  python3 process_photos.py --trip-start 2026-06-06 --trip-end 2026-06-08
                                          # explicit Day-1/last-day dates (e.g. when the folder
                                          # also has photos taken before/after the trip proper);
                                          # overrides auto-detection from photo EXIF dates
  python3 process_photos.py --limit 150 --slideshow
                                          # only process the first 150 source photos in
                                          # chronological order (rest left untouched on disk);
                                          # useful for trying out a large incoming batch
                                          # incrementally instead of all at once
"""

import argparse, io, json, math, os, re, shlex, subprocess, sys, time, urllib.parse, urllib.request
from datetime import datetime, timedelta
from pathlib import Path

from PIL import Image, ImageDraw, ImageOps
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.section import WD_ORIENT
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Configuration ─────────────────────────────────────────────────────────────

FOLDER = Path(__file__).parent

# Trip identity — auto-derived from the containing folder name so this script
# needs zero per-trip editing to reuse in a new folder.
TRIP_NAME = FOLDER.name                                        # display title, e.g. "Sykes"
TRIP_SLUG = re.sub(r'[^A-Za-z0-9]+', '_', TRIP_NAME).strip('_') or 'Trip'   # filesystem-safe
TRIP_TITLE = f'{TRIP_NAME} — Photo Manifest'

# Trip date range — NOT derived from the folder name (folder names rarely carry
# exact start/end days); computed instead from the photos' own EXIF timestamps
# in main() before processing, and used by parse_exif() for "(Day N)" tags.
TRIP_START_DATE = None
TRIP_END_DATE   = None

# ± window around [TRIP_START_DATE, TRIP_END_DATE] before an EXIF DateTimeOriginal
# is flagged 'suspect_date' rather than trusted outright. Generous on purpose —
# people commonly arrive a couple days early or linger after, so this only needs
# to catch genuinely wrong dates (wrong year, camera clock never set), not early/
# late arrivals.
DATE_TRUST_BUFFER_DAYS = 4

STATE_FILE = FOLDER / 'photo_state.json'
MANIFEST   = FOLDER / f'{TRIP_SLUG}_Photo_Manifest.docx'
SLIDES_OUT = FOLDER / f'{TRIP_SLUG}_Photo_Manifest_slides.pptx'
# Auto-detected, not currently consumed elsewhere in the script (reserved for future GPX-track use).
GPX_FILE   = next(FOLDER.glob('*.gpx'), None)
# Scratch dir for orientation-corrected copies of original photos used only in the
# slideshow (ffmpeg's image demuxer ignores JPEG EXIF Orientation — see make_oriented_copy).
# Regenerated on every slideshow build and removed by the encode script when it finishes.
ORIENTED_DIR = FOLDER / '.slideshow_orig_oriented'

GENERATED_SUFFIXES = ('_topo', '_fullinfo', '_nearestarea')

# Any jpg whose stem ends with one of these is a derivative we produced —
# never treat it as a source photo regardless of what else is in the folder.
GENERATED_STEM_ENDINGS = GENERATED_SUFFIXES

FFMPEG            = '/usr/bin/ffmpeg'
SLIDESHOW_OUT     = FOLDER / f'{TRIP_SLUG}_Slideshow.mp4'
TOPO_SKIP_DIST    = 3657.6   # 4000 yards in metres — skip topo if closer than this (fallback, single-photo case)
TOPO_ADAPTIVE_MULT = 3.0     # adaptive threshold = median inter-photo gap × this (doubled from original 1.5)
SLIDESHOW_FPS     = 24
USER_AGENT        = f'{TRIP_SLUG}-photo-captioner/1.0'
# Nominatim's hard cap is 1 req/sec; 1.1s already respects that. Bumped default to 3.0s
# for large one-time bulk runs (hundreds of photos in one sitting) since their policy
# flags "systematic queries" for banning regardless of per-request pacing — see
# [[date_review_workflow]] memory for the underlying discussion. Override with
# --nominatim-delay for a quick small batch where 1.1s is plenty.
NOMINATIM_DELAY_SECONDS = 3.0

# Camera fingerprint table: (make.lower(), model.lower(), software) -> owner
# Multiple keys may map to the same owner (e.g. Robin's old + new phone) — a
# plain dict lookup already handles that with no extra code.
CAMERA_OWNERS = {
    ('samsung', 'galaxy z flip5', 'F731U1UES8GZE8'): 'Bill',
    ('samsung', 'galaxy z flip5', 'F731U1UES7FZC5'): 'Bill',   # older firmware, same phone
    ('apple', 'iphone 17', '26.5.1'):                'Austin',
    ('apple', 'iphone 13', '26.5'):                  'Rob',
    ('samsung', 'galaxy s25+', 'S936USQUACZF1'):     'Frank',
    ('samsung', 'galaxy s25+', 'S936USQSACZE1'):     'Frank',  # older firmware, same phone
    ('apple', 'iphone 14 pro', '26.4.2'):            'Robin',
    ('fujifilm', 'finepix xp140 xp141 xp145',
     'Digital Camera FinePix XP140 XP141 XP145 Ver1.03'): 'Austin',  # rugged/waterproof cam, separate from his iPhone 17
    # Add Steven, and Robin's older phone, when identified
}

# Per-camera fixed EXIF DateTimeOriginal correction, keyed the same as CAMERA_OWNERS.
# Applied in parse_exif() BEFORE the trusted-window check, so a photo whose corrected
# date now falls inside the trip window is treated as fully valid (no date_flag,
# elevation/geocode/topo all run normally) rather than landing on the review list.
# Derive each entry from a single manually-confirmed reference photo, not a guess —
# find one photo from the misbehaving camera whose real date/time you know for certain,
# compute (real datetime - EXIF datetime), and that's the timedelta. Only valid for
# photos taken while the camera's clock held steady at that exact offset; if the clock
# gets reset again, a second entry (or a date-range-scoped correction) will be needed
# instead of one flat delta. Intentionally empty here — these are per-incident facts
# about one specific camera's clock error on one specific trip, NOT reusable across
# trips like CAMERA_OWNERS is, so don't accumulate entries here as a running log.
CAMERA_TIME_OFFSETS = {
    # ('make', 'model', 'software'): timedelta(...),
}

# Manual per-file date/photographer overrides, keyed by filename — for photos with
# NO usable EXIF at all (commonly a photo re-saved by a messaging app, which strips
# Make/Model/Software/DateTimeOriginal/GPS, so there's no camera fingerprint to key a
# CAMERA_TIME_OFFSETS entry off of). Applied in parse_exif() in place of the raw EXIF
# date/photographer lookup — user-confirmed, not inferred, so date_flag is never set
# for one of these regardless of the trip window. Intentionally empty here — these are
# one-off, per-photo facts specific to one trip's folder, not reusable across trips.
MANUAL_OVERRIDES = {
    # 'filename.jpeg': {'datetime': datetime(YYYY, M, D, H, M, S), 'photographer': 'Name'},
}

# ImageMagick caption style
IM_FONT  = 'Helvetica-Bold'
IM_COLOR = 'rgb(255,215,0)'
IMG_W, IMG_H = None, None   # filled per-photo from EXIF

BURN_IN_TOPO = False   # set from --burn-in in main(); topo inset on fullinfo is opt-in

# ── Helpers ───────────────────────────────────────────────────────────────────

def is_source_photo(path: Path) -> bool:
    """True only for original (non-generated) JPEG photos."""
    if path.suffix.lower() not in ('.jpg', '.jpeg'):
        return False
    # Exclude anything whose stem ends with a suffix we generate
    if any(path.stem.endswith(s) for s in GENERATED_STEM_ENDINGS):
        return False
    # Exclude known non-photo jpegs by name pattern (belt-and-suspenders)
    if path.stem.startswith('_'):
        return False
    return True

def file_sig(path: Path) -> tuple:
    st = os.stat(path)
    return (round(st.st_mtime, 2), st.st_size)

def haversine(lat1, lon1, lat2, lon2) -> float:
    R = 6371000
    phi1, phi2 = math.radians(lat1), math.radians(lat2)
    a = (math.sin((phi2 - phi1) / 2) ** 2 +
         math.cos(phi1) * math.cos(phi2) * math.sin(math.radians(lon2 - lon1) / 2) ** 2)
    return R * 2 * math.atan2(math.sqrt(a), math.sqrt(1 - a))

def photo_sort_key(path: Path, state: dict) -> tuple:
    """
    Three-level sort: EXIF datetime (best) → file mtime → filename.
    Returns a tuple that sorts lexicographically in the right order.
    Priority 0 = EXIF datetime, priority 1 = mtime, priority 2 = name only.
    """
    fname = path.name

    # Level 1: EXIF datetime already cached in state
    if fname in state and state[fname].get('datetime_raw'):
        return (0, state[fname]['datetime_raw'], fname)

    # Level 2: read DateTimeOriginal directly from EXIF (fast single-field query)
    try:
        raw = subprocess.check_output(
            ['identify', '-format', '%[EXIF:DateTimeOriginal]', str(path)],
            text=True, stderr=subprocess.DEVNULL, timeout=10).strip()
        if raw and len(raw) >= 19 and ':' in raw:
            return (0, raw, fname)
    except Exception:
        pass

    # Level 3: file modification time expressed as a comparable datetime string
    mtime = os.stat(path).st_mtime
    dt_str = datetime.fromtimestamp(mtime).strftime('%Y:%m:%d %H:%M:%S')
    return (1, dt_str, fname)

def load_state() -> dict:
    if STATE_FILE.exists():
        with open(STATE_FILE) as f:
            return json.load(f)
    return {}

def save_state(state: dict):
    with open(STATE_FILE, 'w') as f:
        json.dump(state, f, indent=2)

# ── EXIF extraction ───────────────────────────────────────────────────────────

def extract_exif(photo: Path) -> dict:
    out = subprocess.check_output(['identify', '-verbose', str(photo)], text=True)
    exif = {}
    for line in out.splitlines():
        line = line.strip()
        if line.startswith('exif:'):
            k, _, v = line[5:].partition(': ')
            exif[k.strip()] = v.strip()
    return exif

def dms_to_dd(dms: str, ref: str) -> float:
    parts = [x.strip() for x in dms.split(',')]
    def frac(s):
        n, _, d = s.partition('/')
        return float(n) / float(d)
    dd = frac(parts[0]) + frac(parts[1]) / 60 + frac(parts[2]) / 3600
    return -dd if ref in ('S', 'W') else dd

def parse_exif(exif: dict, fallback_dt: datetime = None, fname: str = None) -> dict:
    """
    fallback_dt (typically file mtime) stands in when DateTimeOriginal is missing
    or unparseable. date_flag on the return value is:
      None            — EXIF date present and within the trusted trip window
      'no_date'       — no usable DateTimeOriginal at all; fallback_dt was used
      'suspect_date'  — EXIF date (after any CAMERA_TIME_OFFSETS correction) still
                        falls outside the trip window (± DATE_TRUST_BUFFER_DAYS),
                        e.g. a camera clock that was never set and has no known
                        offset yet. Flagged for manual review.
    A camera-specific fixed offset (CAMERA_TIME_OFFSETS) is applied to a parseable
    EXIF date BEFORE the trust-window check, so a photo from a camera with a known,
    constant clock error is corrected and treated as fully valid rather than flagged.
    A MANUAL_OVERRIDES[fname] entry (for photos with no EXIF at all to key a camera
    offset off of) takes priority over everything above and is never flagged.
    """
    try:
        lat = dms_to_dd(exif['GPSLatitude'], exif['GPSLatitudeRef'])
        lon = dms_to_dd(exif['GPSLongitude'], exif['GPSLongitudeRef'])
    except KeyError:
        lat, lon = None, None

    make     = exif.get('Make', '').strip().lower()
    model    = exif.get('Model', '').strip().lower()
    software = exif.get('Software', '').strip()
    camera_key = (make, model, software)

    override = MANUAL_OVERRIDES.get(fname) if fname else None
    photographer = (override.get('photographer') if override and override.get('photographer')
                    else CAMERA_OWNERS.get(camera_key, ''))

    raw_exif_date = exif.get('DateTimeOriginal', '').strip()
    date_flag = None
    dt = None
    if override:
        dt = override['datetime']
    elif raw_exif_date:
        try:
            dt = datetime.strptime(raw_exif_date, '%Y:%m:%d %H:%M:%S')
        except ValueError:
            dt = None
    if dt is None:
        date_flag = 'no_date'
        dt = fallback_dt or datetime.now()
    elif not override:
        offset = CAMERA_TIME_OFFSETS.get(camera_key)
        if offset:
            dt = dt + offset
        if TRIP_START_DATE and TRIP_END_DATE:
            buf = timedelta(days=DATE_TRUST_BUFFER_DAYS)
            if not (TRIP_START_DATE - buf <= dt.date() <= TRIP_END_DATE + buf):
                date_flag = 'suspect_date'

    w   = int(exif.get('ImageWidth', 0))
    h   = int(exif.get('ImageLength', 0))

    day_tag = ''
    if date_flag is None and TRIP_START_DATE and TRIP_END_DATE and TRIP_START_DATE <= dt.date() <= TRIP_END_DATE:
        day_tag = f' (Day {(dt.date() - TRIP_START_DATE).days + 1})'
    elif date_flag:
        day_tag = '  [DATE NEEDS REVIEW]'

    return {
        'lat': lat, 'lon': lon,
        'datetime_raw': dt.strftime('%Y:%m:%d %H:%M:%S'),
        'datetime_str': dt.strftime('%B %-d, %Y  at  %-I:%M %p PDT') + day_tag,
        'img_w': w, 'img_h': h,
        'photographer': photographer,
        'camera_key': (make, model, software),
        'date_flag': date_flag,
        'date_flag_raw_exif': raw_exif_date,
    }

# ── Elevation (USGS) ──────────────────────────────────────────────────────────

def get_elevation(lat: float, lon: float) -> tuple:
    url = (f'https://epqs.nationalmap.gov/v1/json'
           f'?x={lon:.6f}&y={lat:.6f}&units=Feet&wkid=4326&includeDate=false')
    for attempt in range(4):
        try:
            with urllib.request.urlopen(url, timeout=15) as r:
                ft = json.load(r)['value']
            return ft, ft * 0.3048
        except Exception as e:
            print(f'  USGS attempt {attempt+1} failed: {e}')
            time.sleep(2 ** attempt)
    return None, None

# ── Reverse geocode (Nominatim zoom=17 + Overpass natural features) ───────────

NATURAL_SEARCH_RADIUS = 3000   # metres — look for lakes/peaks within this range

def _nearest_natural_feature(lat: float, lon: float) -> tuple[str, float]:
    """
    Query Overpass for the nearest named lake or peak within NATURAL_SEARCH_RADIUS.
    Uses out geom so lake polygons are measured to their nearest vertex (not centroid),
    which correctly identifies large lakes when the photo is taken on or near the shore.
    Returns (name, distance_m) or ('', inf) on failure / nothing found.
    """
    query = (
        f'[out:json][timeout:12];'
        f'('
        f'node(around:{NATURAL_SEARCH_RADIUS},{lat:.6f},{lon:.6f})'
        f'["natural"~"peak|saddle"]["name"];'
        f'way(around:{NATURAL_SEARCH_RADIUS},{lat:.6f},{lon:.6f})'
        f'["natural"~"water|lake"]["name"];'
        f'relation(around:{NATURAL_SEARCH_RADIUS},{lat:.6f},{lon:.6f})'
        f'["natural"~"water|lake"]["name"];'
        f');out geom;'
    )
    url = ('https://overpass-api.de/api/interpreter?data='
           + urllib.parse.quote(query))
    req = urllib.request.Request(url, headers={'User-Agent': USER_AGENT})

    for attempt in range(2):
        try:
            with urllib.request.urlopen(req, timeout=15) as r:
                data = json.load(r)
            break
        except Exception as e:
            if attempt == 0:
                time.sleep(3)
            else:
                print(f'  Overpass natural lookup failed: {e}')
                return '', float('inf')

    best_name, best_dist = '', float('inf')
    for el in data.get('elements', []):
        name = el.get('tags', {}).get('name', '')
        if not name:
            continue
        if el['type'] == 'node':
            pts = [(el['lat'], el['lon'])]
        elif el['type'] == 'way':
            # Nearest vertex on polygon boundary — accurate for large lakes.
            pts = [(p['lat'], p['lon']) for p in el.get('geometry', [])]
        elif el['type'] == 'relation':
            # Multipolygon: gather vertices from all member geometries.
            pts = [
                (p['lat'], p['lon'])
                for m in el.get('members', [])
                for p in m.get('geometry', [])
            ]
        else:
            continue
        if not pts:
            continue
        d = min(haversine(lat, lon, p[0], p[1]) for p in pts)
        if d < best_dist:
            best_dist, best_name = d, name

    return best_name, best_dist

def _nominatim_reverse(lat: float, lon: float, zoom: int) -> dict:
    time.sleep(NOMINATIM_DELAY_SECONDS)
    req = urllib.request.Request(
        f'https://nominatim.openstreetmap.org/reverse'
        f'?format=json&lat={lat:.6f}&lon={lon:.6f}&zoom={zoom}',
        headers={'User-Agent': USER_AGENT})
    with urllib.request.urlopen(req, timeout=15) as r:
        return json.load(r)

def get_area_phrase(lat: float, lon: float) -> str:
    # ── Step 1a: Nominatim zoom=17 ────────────────────────────────────────────
    nom  = _nominatim_reverse(lat, lon, 17)
    addr = nom.get('address', {})

    road    = addr.get('road')
    nat_nom = addr.get('natural') or addr.get('leisure')

    # ── Step 1b: If no natural/leisure at zoom=17, try zoom=18 ───────────────
    # zoom=18 surfaces leisure features (ski areas, parks) that zoom=17 misses
    # when a minor road/trail name dominates at the coarser zoom.
    if not nat_nom:
        nom18  = _nominatim_reverse(lat, lon, 18)
        addr18 = nom18.get('address', {})
        nat_nom = addr18.get('natural') or addr18.get('leisure')

    # natural/leisure is more meaningful than a road name when both are present
    specific = nat_nom or road or addr.get('amenity') or addr.get('hamlet') or addr.get('suburb')
    locality = addr.get('town') or addr.get('village') or addr.get('county')
    nom_phrase = ', '.join(p for p in [specific, locality] if p)

    # ── Step 2: Overpass natural feature supplement ───────────────────────────
    # Only query if Nominatim gave us a trail/road (or nothing specific at all)
    # and skip if Nominatim already found a natural/leisure feature.
    if nat_nom:
        return nom_phrase   # Nominatim already found a natural feature — done

    nat_name, nat_dist = _nearest_natural_feature(lat, lon)

    if not nat_name or nat_dist >= NATURAL_SEARCH_RADIUS:
        return nom_phrase   # nothing useful found nearby

    # Suppress if the natural feature name is already embedded in the road name
    # (e.g. "Horseshoe Lake Parking" + "Horseshoe Lake" → redundant)
    if specific and nat_name.lower() in specific.lower():
        return nom_phrase

    if specific:
        # On a trail/road — append the nearest natural feature
        return f'{specific}, near {nat_name}'
    else:
        # No specific Nominatim result at all — lead with natural feature
        locality_str = locality or ''
        return ', '.join(p for p in [f'near {nat_name}', locality_str] if p)

# ── Topo map tile fetch ───────────────────────────────────────────────────────

def world_px(lat, lon, zoom):
    n = 2 ** zoom
    x = (lon + 180) / 360 * n * 256
    y = (1 - math.asinh(math.tan(math.radians(lat))) / math.pi) / 2 * n * 256
    return x, y

def fetch_topo(lat: float, lon: float, out_path: Path, zoom=16, out_size=1000):
    cx, cy   = world_px(lat, lon, zoom)
    half     = out_size // 2
    tx_min   = int((cx - half) // 256)
    tx_max   = int((cx + half) // 256)
    ty_min   = int((cy - half) // 256)
    ty_max   = int((cy + half) // 256)

    canvas = Image.new('RGB', ((tx_max - tx_min + 1) * 256,
                                (ty_max - ty_min + 1) * 256))
    print(f'  Fetching {(tx_max-tx_min+1)*(ty_max-ty_min+1)} topo tiles...')

    for ty in range(ty_min, ty_max + 1):
        for tx in range(tx_min, tx_max + 1):
            url = (f'https://basemap.nationalmap.gov/arcgis/rest/services/'
                   f'USGSTopo/MapServer/tile/{zoom}/{ty}/{tx}')
            req = urllib.request.Request(url, headers={'User-Agent': USER_AGENT})
            for attempt in range(3):
                try:
                    with urllib.request.urlopen(req, timeout=15) as r:
                        tile = Image.open(io.BytesIO(r.read())).convert('RGB')
                    break
                except Exception as e:
                    print(f'    tile {tx},{ty} attempt {attempt+1}: {e}')
                    time.sleep(1)
            canvas.paste(tile, ((tx - tx_min) * 256, (ty - ty_min) * 256))
            time.sleep(0.05)

    crop_x = int(cx - tx_min * 256) - half
    crop_y = int(cy - ty_min * 256) - half
    cropped = canvas.crop((crop_x, crop_y, crop_x + out_size, crop_y + out_size))

    draw = ImageDraw.Draw(cropped)
    r = 12
    draw.line([(half - r, half), (half + r, half)], fill='red', width=2)
    draw.line([(half, half - r), (half, half + r)], fill='red', width=2)
    draw.ellipse([(half - 5, half - 5), (half + 5, half + 5)], outline='red', width=2)

    cropped.save(str(out_path), 'JPEG', quality=90)

# ── ImageMagick caption images ────────────────────────────────────────────────

def _caption_lines(meta: dict):
    """Return caption lines from metadata. Omits GPS/area lines if no coordinates."""
    lines = [meta['datetime_str']]
    if meta.get('lat') is not None:
        if meta.get('elev_ft') is not None:
            elev_str = (f"{meta['elev_ft']:.0f} ft  ({meta['elev_m']:.0f} m),   "
                        f"{meta['lat']:.5f}°N,  {abs(meta['lon']):.5f}°W")
        else:
            elev_str = f"{meta['lat']:.5f}°N,  {abs(meta['lon']):.5f}°W"
        lines.append(elev_str)
        if meta.get('area_phrase'):
            lines.append(meta['area_phrase'])
    return lines

def make_fullinfo(photo: Path, topo: Path, dst: Path, meta: dict):
    # Always output 4K letterboxed so the topo inset and caption never
    # overlap the photo content on portrait or non-16:9 shots.
    W, H = 3840, 2160
    pt   = 56
    lh   = int(pt * 1.55)
    pad  = int(pt * 0.9)

    lines   = _caption_lines(meta)
    strip_h = len(lines) * lh + 2 * pad

    annot_args = []
    for i, line in enumerate(reversed(lines)):
        annot_args += ['-annotate', f'+0+{pad + i*lh}', line]

    topo_size = int(H * 0.20 * 2.5)
    mx = int(W * 0.025)
    my = int(H * 0.025)

    cmd = (
        ['convert', str(photo),
         '-auto-orient',          # apply EXIF rotation physically before any resize
         '-resize', f'{W}x{H}',
         '-background', 'black',
         '-gravity', 'Center',
         '-extent', f'{W}x{H}'] +
        ['(', '-size', f'{W}x{strip_h}', 'xc:rgba(0,0,0,0.30)', ')'] +
        ['-gravity', 'South', '-composite'] +
        ['-font', IM_FONT, '-pointsize', str(pt), '-fill', IM_COLOR, '-gravity', 'South'] +
        annot_args +
        ((['(', str(topo), '-resize', f'{topo_size}x{topo_size}', ')'] +
          ['-gravity', 'SouthEast', '-geometry', f'+{mx}+{my}', '-composite'])
         if BURN_IN_TOPO and topo and topo.exists() else []) +
        (['-font', IM_FONT, '-pointsize', str(pt), '-fill', IM_COLOR,
          '-gravity', 'SouthWest', '-annotate', f'+{mx}+{my}',
          f'({meta["photographer"]})']
         if meta['photographer'] else []) +
        ['-quality', '95', str(dst)]
    )

    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f'  ImageMagick error: {result.stderr[:200]}')

def make_oriented_copy(photo: Path, dst: Path):
    """Physically apply EXIF rotation to a copy of the original photo.

    ffmpeg's image demuxer ignores JPEG EXIF Orientation (unlike ImageMagick,
    which -auto-orient already handles for fullinfo/nearestarea/topo), so the
    "original" slide in the slideshow needs its own pre-rotated copy or it
    plays back rotated relative to the corrected slides around it."""
    cmd = ['convert', str(photo), '-auto-orient', '-quality', '95', str(dst)]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f'  ImageMagick error (orient original): {result.stderr[:200]}')

def make_nearestarea(photo: Path, dst: Path, meta: dict):
    W, H = 3840, 2160
    pt    = 56
    h_pad = int(pt * 0.75)   # horizontal clearance on each side of text
    v_pad = int(pt * 0.45)   # vertical clearance above and below text
    bot   = int(pt * 0.35)   # gap from image bottom to bottom of box

    # Measure the rendered text width/height using ImageMagick label pseudo-image
    m = subprocess.run(
        ['convert', '-font', IM_FONT, '-pointsize', str(pt),
         f'label:{meta["area_phrase"]}', '-format', '%wx%h', 'info:'],
        capture_output=True, text=True
    )
    try:
        txt_w, txt_h = map(int, m.stdout.strip().split('x'))
    except ValueError:
        txt_w, txt_h = int(len(meta['area_phrase']) * pt * 0.55), int(pt * 1.2)

    box_w = min(txt_w + 2 * h_pad, W)
    box_h = txt_h + 2 * v_pad
    bx0   = (W - box_w) // 2
    bx1   = bx0 + box_w
    by1   = H - bot
    by0   = by1 - box_h
    # With -gravity South -annotate +0+ann_y the baseline is ann_y px from bottom.
    # Setting ann_y = bot + v_pad puts the baseline v_pad above the box bottom,
    # giving equal padding above and below the text inside the box.
    ann_y = bot + v_pad

    cmd = (
        ['convert', str(photo),
         '-auto-orient',
         '-resize', f'{W}x{H}',
         '-background', 'black',
         '-gravity', 'Center',
         '-extent', f'{W}x{H}',
         '-fill', 'rgba(0,0,0,0.45)',
         '-draw', f'roundrectangle {bx0},{by0} {bx1},{by1} 18,18',
         '-font', IM_FONT, '-pointsize', str(pt), '-fill', IM_COLOR,
         '-gravity', 'South', '-annotate', f'+0+{ann_y}', meta['area_phrase'],
         '-quality', '95', str(dst)]
    )
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f'  ImageMagick error: {result.stderr[:200]}')

def make_topo_captioned(topo: Path, meta: dict):
    """Burn the same 3-line fullinfo caption onto the topo image in-place.
    Call this AFTER make_fullinfo so the inset in fullinfo uses the clean topo."""
    W, H = 1000, 1000
    # Scale pt proportionally to topo height so text fits the smaller canvas.
    pt  = int(56 * H / 2160)
    lh  = int(pt * 1.55)
    pad = int(pt * 0.9)

    lines   = _caption_lines(meta)
    strip_h = len(lines) * lh + 2 * pad

    annot_args = []
    for i, line in enumerate(reversed(lines)):
        annot_args += ['-annotate', f'+0+{pad + i*lh}', line]

    mx = int(W * 0.025)
    my = int(H * 0.025)

    cmd = (
        ['convert', str(topo)] +
        ['(', '-size', f'{W}x{strip_h}', 'xc:rgba(0,0,0,0.30)', ')'] +
        ['-gravity', 'South', '-composite'] +
        ['-font', IM_FONT, '-pointsize', str(pt), '-fill', IM_COLOR, '-gravity', 'South'] +
        annot_args +
        (['-font', IM_FONT, '-pointsize', str(pt), '-fill', IM_COLOR,
          '-gravity', 'SouthWest', '-annotate', f'+{mx}+{my}',
          f'({meta["photographer"]})']
         if meta['photographer'] else []) +
        ['-quality', '90', str(topo)]
    )

    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f'  ImageMagick error (topo caption): {result.stderr[:200]}')

# ── Date review report ─────────────────────────────────────────────────────────

DATE_REVIEW_REPORT = FOLDER / f'{TRIP_SLUG}_DATE_REVIEW_NEEDED.txt'

def write_date_review_report(state: dict) -> int:
    """
    Text list of every photo currently flagged 'no_date' or 'suspect_date' in state
    (see parse_exif). Regenerated from state on every run so it stays in sync even
    across incremental (non --force) runs. Deleted if nothing is flagged.
    """
    flagged = sorted(
        ((fn, meta) for fn, meta in state.items() if meta.get('date_flag')),
        key=lambda kv: kv[0]
    )
    if not flagged:
        if DATE_REVIEW_REPORT.exists():
            DATE_REVIEW_REPORT.unlink()
        return 0

    window = (f'{TRIP_START_DATE} .. {TRIP_END_DATE}'
              if TRIP_START_DATE else 'not yet established')
    lines = [
        f'{len(flagged)} photo(s) need a manual date/time review',
        f'generated {datetime.now().strftime("%Y-%m-%d %H:%M")}  —  trusted trip window: {window}',
        '=' * 78,
    ]
    for fname, meta in flagged:
        flag = meta['date_flag']
        who  = meta.get('photographer') or '(unattributed camera)'
        if flag == 'no_date':
            detail = 'no DateTimeOriginal in EXIF at all — file mtime used as placeholder'
        else:
            detail = f'EXIF says {meta.get("date_flag_raw_exif")} — outside the trusted trip window'
        lines.append(f'{fname:34s} [{flag:13s}] {who:14s} {detail}')

    DATE_REVIEW_REPORT.write_text('\n'.join(lines) + '\n')
    return len(flagged)

# ── Recaption (no elevation/geocode network calls) ────────────────────────────

def recaption_all(state: dict):
    """
    Regenerate every already-processed photo's _fullinfo/_nearestarea/topo-caption
    using each entry's already-cached lat/lon/elev_ft/elev_m/area_phrase — reused
    as-is, no USGS elevation or Nominatim/Overpass calls. Only datetime_str/date_flag
    are freshly recomputed (against the current TRIP_START_DATE/TRIP_END_DATE and
    CAMERA_TIME_OFFSETS), which is all that's needed after changing the trip date
    range. Topo tiles ARE re-fetched from USGS's basemap tile server (permissive,
    no rate-limit policy) because make_topo_captioned burns the caption into that
    file in-place — the clean base doesn't survive the first captioning.
    A photo that newly resolves out of a flagged state gets fully generated; one
    that newly becomes flagged (e.g. a narrower trip window with a photo now
    outside DATE_TRUST_BUFFER_DAYS) has its derived files removed and drops back
    to review-list-only, matching normal processing behavior.
    """
    print(f'\nRecaptioning {len(state)} photos against trip window '
          f'{TRIP_START_DATE} .. {TRIP_END_DATE} (no elevation/geocode calls)...')
    for fname, cached in sorted(state.items()):
        photo = FOLDER / fname
        if not photo.exists():
            print(f'  Skipping {fname} — file missing')
            continue
        stem = photo.stem

        exif_raw    = extract_exif(photo)
        fallback_dt = datetime.fromtimestamp(photo.stat().st_mtime)
        parsed      = parse_exif(exif_raw, fallback_dt=fallback_dt, fname=fname)

        if parsed['date_flag']:
            print(f'  {fname}: now flagged ({parsed["date_flag"]}) — '
                  f'removing derived files, review-list only')
            for suffix in ('_topo.jpg', '_fullinfo.jpg', '_nearestarea.jpg'):
                p = FOLDER / f'{stem}{suffix}'
                if p.exists():
                    p.unlink()
            state[fname] = {
                'sig':                cached.get('sig'),
                'date_flag':          parsed['date_flag'],
                'date_flag_raw_exif': parsed['date_flag_raw_exif'],
                'datetime_raw':       parsed['datetime_raw'],
                'photographer':       parsed['photographer'],
            }
            save_state(state)
            continue

        meta = {
            **parsed,
            'elev_ft':     cached.get('elev_ft'),
            'elev_m':      cached.get('elev_m'),
            'area_phrase': cached.get('area_phrase', ''),
        }
        has_gps = meta['lat'] is not None
        print(f'  {fname}: {meta["datetime_str"]}')

        topo_fname = topo_path = None
        if has_gps:
            topo_fname = f'{stem}_topo.jpg'
            topo_path  = FOLDER / topo_fname
            fetch_topo(meta['lat'], meta['lon'], topo_path)

        fi_fname = f'{stem}_fullinfo.jpg'
        na_fname = f'{stem}_nearestarea.jpg' if has_gps else None
        make_fullinfo(photo, topo_path, FOLDER / fi_fname, meta)
        if has_gps:
            make_nearestarea(photo, FOLDER / na_fname, meta)
            make_topo_captioned(topo_path, meta)

        state[fname] = {
            'sig':                cached.get('sig'),
            'datetime_raw':       parsed['datetime_raw'],
            'datetime_str':       parsed['datetime_str'],
            'lat':                meta['lat'], 'lon': meta['lon'],
            'img_w':              parsed['img_w'], 'img_h': parsed['img_h'],
            'elev_ft':            meta['elev_ft'], 'elev_m': meta['elev_m'],
            'area_phrase':        meta['area_phrase'],
            'photographer':       parsed['photographer'],
            'topo':               topo_fname,
            'fullinfo':           fi_fname,
            'nearestarea':        na_fname,
            'date_flag':          None,
            'date_flag_raw_exif': parsed['date_flag_raw_exif'],
        }
        save_state(state)
    print('Recaption pass complete.')

# ── Manifest builder ──────────────────────────────────────────────────────────

def _thumbnail_bytes(path: Path, max_dim: int = 900, quality: int = 85) -> io.BytesIO:
    """
    Resized/recompressed JPEG copy of an image, in-memory, for embedding in the
    manifest. add_picture() on the original full-resolution source file embeds its
    actual bytes regardless of the tiny on-page display height — across ~800 photos
    that bloated the manifest to gigabytes for no visible benefit at 1.5" tall. This
    caps the long edge at max_dim (900px comfortably covers print-quality at the
    manifest's display size) and re-encodes as JPEG, cutting file size by ~30-40x.
    exif_transpose bakes in EXIF rotation physically, since the resize discards the
    orientation tag otherwise.
    """
    img = ImageOps.exif_transpose(Image.open(path))
    if img.mode != 'RGB':
        img = img.convert('RGB')
    img.thumbnail((max_dim, max_dim), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, 'JPEG', quality=quality)
    buf.seek(0)
    return buf

def build_manifest(state: dict):
    print('\nBuilding manifest...')
    doc = Document()
    section = doc.sections[0]
    section.orientation   = WD_ORIENT.LANDSCAPE
    section.page_width    = Inches(11)
    section.page_height   = Inches(8.5)
    section.top_margin    = Inches(0.5)
    section.bottom_margin = Inches(0.5)
    section.left_margin   = Inches(0.5)
    section.right_margin  = Inches(0.5)

    title_para = doc.add_paragraph()
    run = title_para.add_run(TRIP_TITLE)
    run.bold = True
    run.font.size = Pt(16)
    title_para.paragraph_format.space_after = Pt(8)

    COL_PHOTO = Inches(2.75)
    COL_TOPO  = Inches(1.60)
    COL_INFO  = Inches(10 - 2.75 - 1.60)
    IMG_H     = Inches(1.5)

    def set_col_width(cell, width):
        tcPr = cell._tc.get_or_add_tcPr()
        tcW  = OxmlElement('w:tcW')
        tcW.set(qn('w:w'), str(int(width.inches * 1440)))
        tcW.set(qn('w:type'), 'dxa')
        tcPr.append(tcW)

    table = doc.add_table(rows=1, cols=3)
    table.style = 'Table Grid'

    # Header
    hdr = table.rows[0].cells
    for cell, label, w in zip(hdr, ['Photo', 'Info', 'Topo'],
                               [COL_PHOTO, COL_INFO, COL_TOPO]):
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(label)
        run.bold = True
        run.font.size = Pt(10)
        set_col_width(cell, w)
        cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER

    # Sort photos chronologically
    # Photos still flagged (no usable date yet) are excluded — they have no
    # derived topo/fullinfo/nearestarea files to show, and no confirmed place in
    # the chronological order. They only ever appear in DATE_REVIEW_NEEDED.txt.
    entries = sorted(
        (kv for kv in state.items() if not kv[1].get('date_flag')),
        key=lambda kv: kv[1].get('datetime_raw', '')
    )

    for fname, meta in entries:
        photo_path = FOLDER / fname
        topo_path  = FOLDER / meta['topo'] if meta.get('topo') else None

        if not photo_path.exists():
            print(f'  Skipping {fname} — file missing')
            continue

        row   = table.add_row()
        cells = row.cells
        for cell, w in zip(cells, [COL_PHOTO, COL_INFO, COL_TOPO]):
            set_col_width(cell, w)

        # Photo thumbnail
        p0 = cells[0].paragraphs[0]
        p0.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p0.add_run().add_picture(_thumbnail_bytes(photo_path), height=IMG_H)
        cells[0].vertical_alignment = WD_ALIGN_VERTICAL.CENTER

        # Info
        info_rows = [('Date/Time', meta['datetime_str'])]
        if meta.get('lat') is not None:
            elev_val = (f"{meta['elev_ft']:.0f} ft  ({meta['elev_m']:.0f} m)"
                        if meta.get('elev_ft') is not None else 'N/A')
            info_rows += [
                ('Elevation', elev_val),
                ('GPS',       f"{meta['lat']:.5f}°N,  {abs(meta['lon']):.5f}°W"),
                ('Area',      meta.get('area_phrase', '')),
            ]
        if meta.get('photographer'):
            info_rows.append(('Camera', meta['photographer']))

        cells[1].paragraphs[0].clear()
        first = True
        for label, value in info_rows:
            p = cells[1].paragraphs[0] if first else cells[1].add_paragraph()
            first = False
            p.paragraph_format.space_before = Pt(2)
            p.paragraph_format.space_after  = Pt(2)
            lb = p.add_run(f'{label}:  ')
            lb.bold = True
            lb.font.size = Pt(10)
            vr = p.add_run(value)
            vr.font.size = Pt(10)
        cells[1].vertical_alignment = WD_ALIGN_VERTICAL.CENTER

        # Topo
        p2 = cells[2].paragraphs[0]
        p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if topo_path and topo_path.exists():
            p2.add_run().add_picture(_thumbnail_bytes(topo_path, max_dim=500), height=IMG_H)
        else:
            p2.add_run('(no GPS)')
        cells[2].vertical_alignment = WD_ALIGN_VERTICAL.CENTER

    doc.save(str(MANIFEST))
    print(f'  Saved: {MANIFEST}  ({os.path.getsize(MANIFEST)//1024} KB)')

# ── PowerPoint builder ────────────────────────────────────────────────────────

def _pptx_black_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

def _pptx_add_fullbleed_image(slide, img_path: Path, slide_w, slide_h):
    """Center the image on the slide, scaled to fit without cropping."""
    with Image.open(img_path) as im:
        iw, ih = im.size
    scale = min(slide_w / iw, slide_h / ih)
    w, h = iw * scale, ih * scale
    left, top = (slide_w - w) / 2, (slide_h - h) / 2
    slide.shapes.add_picture(str(img_path), int(left), int(top), width=int(w), height=int(h))

def build_powerpoint(state: dict):
    print('\nBuilding PowerPoint...')
    prs = Presentation()
    prs.slide_width  = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    blank_layout = prs.slide_layouts[6]   # blank layout

    # Title slide — same title as the Manifest, black background
    title_slide = prs.slides.add_slide(blank_layout)
    _pptx_black_bg(title_slide)
    tb = title_slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(3.15),
        prs.slide_width - PptxInches(1), PptxInches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = TRIP_TITLE
    run.font.bold = True
    run.font.size = PptxPt(40)
    run.font.color.rgb = RGBColor(255, 255, 255)

    # Two slides per photo: topo (if it exists), then fullinfo
    # (fullinfo already reflects whatever --burn-in setting was used to generate it).
    # Photos still flagged (no usable date yet) are excluded — they have no
    # derived topo/fullinfo/nearestarea files to show, and no confirmed place in
    # the chronological order. They only ever appear in DATE_REVIEW_NEEDED.txt.
    entries = sorted(
        (kv for kv in state.items() if not kv[1].get('date_flag')),
        key=lambda kv: kv[1].get('datetime_raw', '')
    )
    for fname, meta in entries:
        topo_path = FOLDER / meta['topo'] if meta.get('topo') else None
        fi_path   = FOLDER / meta['fullinfo'] if meta.get('fullinfo') else None

        if topo_path and topo_path.exists():
            slide = prs.slides.add_slide(blank_layout)
            _pptx_black_bg(slide)
            _pptx_add_fullbleed_image(slide, topo_path, prs.slide_width, prs.slide_height)

        if fi_path and fi_path.exists():
            slide = prs.slides.add_slide(blank_layout)
            _pptx_black_bg(slide)
            _pptx_add_fullbleed_image(slide, fi_path, prs.slide_width, prs.slide_height)

    prs.save(str(SLIDES_OUT))
    print(f'  Saved: {SLIDES_OUT}  ({os.path.getsize(SLIDES_OUT)//1024} KB)')

# ── Slideshow builder ─────────────────────────────────────────────────────────

def build_slideshow(state: dict):
    if not state:
        print('No photos in state — nothing to encode.')
        return

    # Photos still flagged (no usable date yet) are excluded — they have no
    # derived topo/fullinfo/nearestarea files to show, and no confirmed place in
    # the chronological order. They only ever appear in DATE_REVIEW_NEEDED.txt.
    entries = sorted(
        (kv for kv in state.items() if not kv[1].get('date_flag')),
        key=lambda kv: kv[1].get('datetime_raw', '')
    )

    concat_file = FOLDER / '_slideshow_concat.txt'
    lines = ['ffconcat version 1.0']

    ORIENTED_DIR.mkdir(exist_ok=True)

    # ── Adaptive topo threshold: median inter-photo gap × TOPO_ADAPTIVE_MULT ──
    gaps = []
    for i in range(1, len(entries)):
        m0 = entries[i-1][1]
        m1 = entries[i][1]
        if m0.get('lat') is not None and m1.get('lat') is not None:
            gaps.append(haversine(m0['lat'], m0['lon'], m1['lat'], m1['lon']))
    if gaps:
        median_gap = sorted(gaps)[len(gaps) // 2]
        topo_threshold = median_gap * TOPO_ADAPTIVE_MULT
        print(f'  Topo threshold: {topo_threshold/0.9144:.0f} yd  '
              f'(median gap {median_gap/0.9144:.0f} yd × {TOPO_ADAPTIVE_MULT})')
    else:
        topo_threshold = TOPO_SKIP_DIST   # fallback: no GPS gaps to measure
        print(f'  Topo threshold: {topo_threshold/0.9144:.0f} yd  (fallback, no GPS gaps to measure)')

    DUR_TOPO    = 4 * 0.8        # 3.2 s
    DUR_ORIG    = 2 + 1.5        # 3.5 s
    DUR_FULL    = 3
    DUR_NEAREST = 2.5

    last_photo_lat, last_photo_lon = None, None
    topo_count = skipped_topos = 0
    total_secs = 0

    for fname, meta in entries:
        lat, lon = meta.get('lat'), meta.get('lon')

        # ── Topo slide: show only when gap from previous photo exceeds threshold
        include_topo = True
        if lat is None:
            include_topo = False   # no GPS → no topo
        elif last_photo_lat is not None:
            dist_m = haversine(last_photo_lat, last_photo_lon, lat, lon)
            if dist_m < topo_threshold:
                include_topo = False
                skipped_topos += 1
                print(f'  {fname}: topo skipped ({dist_m/0.9144:.0f} yd < {topo_threshold/0.9144:.0f} yd threshold)')

        if include_topo and meta.get('topo'):
            topo_path = FOLDER / meta['topo']
            if topo_path.exists():
                lines += [f"file '{topo_path}'", f'duration {DUR_TOPO}']
                total_secs += DUR_TOPO
                topo_count += 1

        # Always update so next photo compares against this one (only if has GPS)
        if lat is not None:
            last_photo_lat, last_photo_lon = lat, lon

        # ── fullinfo, nearestarea, original ───────────────────────────────────
        for key, dur in [('fullinfo', DUR_FULL), ('nearestarea', DUR_NEAREST)]:
            if not meta.get(key):
                continue
            p = FOLDER / meta[key]
            if p.exists():
                lines += [f"file '{p}'", f'duration {dur}']
                total_secs += dur

        orig = FOLDER / fname
        if orig.exists():
            oriented = ORIENTED_DIR / fname
            make_oriented_copy(orig, oriented)
            lines += [f"file '{oriented}'", f'duration {DUR_ORIG}']
            total_secs += DUR_ORIG

    # Repeat the last file entry without duration so ffmpeg flushes the final frame
    last_file_line = next((l for l in reversed(lines) if l.startswith('file ')), None)
    if last_file_line:
        lines.append(last_file_line)

    with open(concat_file, 'w') as f:
        f.write('\n'.join(lines) + '\n')

    mins, secs = divmod(total_secs, 60)
    print(f'\nSlideshow: {len(entries)} photos, '
          f'{topo_count} topo slides ({skipped_topos} skipped), '
          f'~{mins}m {secs}s total runtime')
    print(f'Output: {SLIDESHOW_OUT}')
    print('Encoding 4K H.264 CRF=14 (this may take several minutes)...\n')

    vf = (
        'scale=3840:2160:force_original_aspect_ratio=decrease,'
        'pad=3840:2160:(ow-iw)/2:(oh-ih)/2:color=black,'
        'setsar=1,format=yuv420p'
    )

    tmp_out = SLIDESHOW_OUT.with_name(SLIDESHOW_OUT.stem + '_novid.mp4')

    cmd_video = [
        FFMPEG, '-y',
        '-f', 'concat', '-safe', '0', '-i', str(concat_file),
        '-vf', vf,
        '-r', str(SLIDESHOW_FPS),
        '-c:v', 'libx264',
        '-crf', '14',
        '-preset', 'fast',
        '-profile:v', 'high',
        '-level:v', '5.1',
        '-an',
        str(tmp_out),
    ]

    cmd_audio = [
        FFMPEG, '-y',
        '-i', str(tmp_out),
        '-f', 'lavfi', '-i', 'anullsrc=r=48000:cl=stereo',
        '-c:v', 'copy',
        '-c:a', 'aac', '-b:a', '32k',
        '-shortest',
        '-movflags', '+faststart',
        str(SLIDESHOW_OUT),
    ]

    # ── Launch encode in an xfce4-terminal so progress is visible ────────────
    encode_script = FOLDER / '_encode_slideshow.sh'
    cmd_video_str = ' '.join(shlex.quote(str(a)) for a in cmd_video)
    cmd_audio_str = ' '.join(shlex.quote(str(a)) for a in cmd_audio)
    bash_safe_trip_name = TRIP_NAME.replace('"', '\\"')   # embedded in a double-quoted bash echo below
    script_body = f"""#!/bin/bash
echo "═══════════════════════════════════════════════════"
echo "  {bash_safe_trip_name} Slideshow — FFmpeg 4K H.264 CRF=14 Encode"
echo "  Output: {SLIDESHOW_OUT}"
echo "═══════════════════════════════════════════════════"
echo ""
echo "Step 1/2 — Video encode..."
echo ""
{cmd_video_str}
STATUS=$?
echo ""
if [ $STATUS -ne 0 ]; then
    echo "  FFmpeg video encode failed (exit $STATUS)"
    rm -rf {shlex.quote(str(ORIENTED_DIR))}
    read -p "Press Enter to close..."
    exit 1
fi

echo "Step 2/2 — Muxing silent 32kbps AAC audio track..."
echo ""
{cmd_audio_str}
STATUS=$?
rm -f {shlex.quote(str(tmp_out))}
rm -rf {shlex.quote(str(ORIENTED_DIR))}
echo ""
if [ $STATUS -eq 0 ]; then
    SIZE=$(du -sh {shlex.quote(str(SLIDESHOW_OUT))} 2>/dev/null | cut -f1)
    echo "═══════════════════════════════════════════════════"
    echo "  Done.  File size: $SIZE"
    echo "═══════════════════════════════════════════════════"
else
    echo "  FFmpeg audio mux failed (exit $STATUS)"
fi
echo ""
read -p "Press Enter to close this window..."
rm -f {shlex.quote(str(encode_script))}
"""
    encode_script.write_text(script_body)
    encode_script.chmod(0o755)

    subprocess.Popen([
        'xfce4-terminal',
        f'--title={TRIP_NAME} Slideshow — FFmpeg Encode',
        f'--geometry=80x25',
        '-e', f'bash {shlex.quote(str(encode_script))}',
    ])
    print(f'\nEncode launched in xfce4-terminal.')
    print(f'Output → {SLIDESHOW_OUT}')
    print('(The terminal will stay open when finished — press Enter there to close it.)')


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    global BURN_IN_TOPO, TRIP_START_DATE, TRIP_END_DATE, NOMINATIM_DELAY_SECONDS

    parser = argparse.ArgumentParser()
    parser.add_argument('--force',        action='store_true',
                        help='Reprocess all photos regardless of state')
    parser.add_argument('--manifest',     action='store_true',
                        help='Rebuild manifest from cache without reprocessing images')
    parser.add_argument('--slideshow',    action='store_true',
                        help='Build slideshow without prompting')
    parser.add_argument('--no-slideshow', action='store_true',
                        help='Skip slideshow prompt and do not encode')
    parser.add_argument('--burn-in',      action='store_true',
                        help='Burn the topo map inset into _fullinfo images (default off). '
                             'Run with --force to regenerate already-processed fullinfo images.')
    parser.add_argument('--powerpoint',   action='store_true',
                        help='Also build a black-background .pptx slide deck '
                             '(title slide + topo/fullinfo slides per photo). Default off.')
    parser.add_argument('--trip-start',   metavar='YYYY-MM-DD',
                        help='Explicit Day-1 date for "(Day N)" tags (overrides auto-detection '
                             'from photo EXIF dates). Photos outside [--trip-start, --trip-end] '
                             'get no Day tag. Must be given together with --trip-end.')
    parser.add_argument('--trip-end',     metavar='YYYY-MM-DD',
                        help='Explicit last-day date for "(Day N)" tags. Must be given together '
                             'with --trip-start.')
    parser.add_argument('--nominatim-delay', type=float, metavar='SECONDS',
                        default=NOMINATIM_DELAY_SECONDS,
                        help=f'Seconds to sleep before each Nominatim reverse-geocode call '
                             f'(default {NOMINATIM_DELAY_SECONDS}; their hard cap is 1.0). '
                             f'Higher is safer for large one-time bulk runs.')
    parser.add_argument('--limit',        type=int, metavar='N',
                        help='Only process the first N source photos in chronological order '
                             '(the rest stay on disk untouched, unprocessed, and out of the '
                             'manifest/slideshow until a later run without --limit, or with a '
                             'higher one, picks them up). TRIP_START_DATE/TRIP_END_DATE are '
                             'still computed from all on-disk photos, not just the first N, so '
                             '"(Day N)" tags stay correct as more batches get processed later.')
    parser.add_argument('--recaption',    action='store_true',
                        help='Regenerate every already-processed photo\'s captions '
                             '(_fullinfo/_nearestarea/topo caption) using each entry\'s '
                             'already-cached lat/lon/elevation/area-phrase — no elevation '
                             'or reverse-geocode network calls, only topo tiles are '
                             're-fetched (caption is burned into that file in-place). Use '
                             'after changing --trip-start/--trip-end so "(Day N)" tags/'
                             'captions update without re-hitting USGS elevation or Nominatim.')
    args = parser.parse_args()

    if args.limit is not None and args.limit <= 0:
        parser.error('--limit must be a positive integer')

    # back-compat alias
    if hasattr(args, 'manifest_only'):
        args.manifest = args.manifest or args.manifest_only

    BURN_IN_TOPO = args.burn_in
    NOMINATIM_DELAY_SECONDS = args.nominatim_delay

    if args.trip_start or args.trip_end:
        if not (args.trip_start and args.trip_end):
            parser.error('--trip-start and --trip-end must be given together')
        try:
            TRIP_START_DATE = datetime.strptime(args.trip_start, '%Y-%m-%d').date()
            TRIP_END_DATE   = datetime.strptime(args.trip_end, '%Y-%m-%d').date()
        except ValueError:
            parser.error('--trip-start/--trip-end must be in YYYY-MM-DD format')
        if TRIP_START_DATE > TRIP_END_DATE:
            parser.error('--trip-start must not be after --trip-end')

    state = load_state()

    if args.recaption:
        recaption_all(state)
        save_state(state)
    elif not args.manifest:
        # ── Detect source photos on disk ──────────────────────────────────────
        on_disk = {p.name for p in FOLDER.iterdir()
                   if p.suffix.lower() in ('.jpg', '.jpeg') and is_source_photo(p)}
        in_state = set(state.keys())

        # Remove entries for deleted photos
        deleted = in_state - on_disk
        for fname in deleted:
            print(f'Removed from state (deleted): {fname}')
            del state[fname]

        # Build time-ordered list: EXIF datetime → file mtime → filename
        source_paths = sorted(
            (FOLDER / f for f in on_disk),
            key=lambda p: photo_sort_key(p, state)
        )
        print(f'\nProcessing order ({len(source_paths)} source photos):')
        photo_dates = []
        # GPS-anchored subset of photo_dates, used (in preference to photo_dates) to
        # compute TRIP_START_DATE/TRIP_END_DATE below. A photo's EXIF date is trusted
        # for this purpose only if it's GPS-tagged — a camera with a bad/unset clock
        # (no GPS) can't silently drag the whole trip's "(Day N)" tags off by a year.
        trusted_dates = []
        for p in source_paths:
            key = photo_sort_key(p, state)
            label = 'EXIF' if key[0] == 0 else 'mtime'
            print(f'  [{label}] {key[1]}  {p.name}')
            d = datetime.strptime(key[1], '%Y:%m:%d %H:%M:%S').date()
            photo_dates.append(d)
            if key[0] == 0:
                cached_lat = state.get(p.name, {}).get('lat')
                if cached_lat is not None:
                    trusted_dates.append(d)
                else:
                    try:
                        gps = subprocess.check_output(
                            ['identify', '-format', '%[EXIF:GPSLatitude]', str(p)],
                            text=True, stderr=subprocess.DEVNULL, timeout=10).strip()
                        if gps:
                            trusted_dates.append(d)
                    except Exception:
                        pass
        print()

        # Trip date range for "(Day N)" tags — derived from the photos themselves,
        # not the folder name, since folder names rarely carry exact start/end days.
        # Skipped if --trip-start/--trip-end were given explicitly above (e.g. when the
        # folder also holds photos taken before/after the trip proper).
        if photo_dates and TRIP_START_DATE is None:
            anchor = trusted_dates or photo_dates
            TRIP_START_DATE, TRIP_END_DATE = min(anchor), max(anchor)

        # --limit trims the processing loop only, applied after the full-corpus scan
        # above so TRIP_START_DATE/TRIP_END_DATE (and thus "(Day N)" tags) reflect every
        # on-disk photo even when just a chronological prefix is actually processed.
        if args.limit is not None and args.limit < len(source_paths):
            print(f'--limit {args.limit}: processing the first {args.limit} of '
                  f'{len(source_paths)} photos (chronological order); the rest stay '
                  f'untouched until a later run.\n')
            source_paths = source_paths[:args.limit]

        # Process in chronological order
        for photo in source_paths:
            fname = photo.name
            sig   = file_sig(photo)
            stem  = photo.stem

            cached = state.get(fname, {})
            sig_matches = not args.force and cached.get('sig') == list(sig)

            if sig_matches and not cached.get('date_flag'):
                # Regenerate any missing derived files from cached state (no API calls).
                stem = photo.stem
                missing = []
                for key, suffix in [('topo', '_topo.jpg'), ('fullinfo', '_fullinfo.jpg'),
                                     ('nearestarea', '_nearestarea.jpg')]:
                    if cached.get(key) and not (FOLDER / f'{stem}{suffix}').exists():
                        missing.append((key, suffix))
                if missing:
                    print(f'Restoring derived files for: {fname}')
                    cmeta = {
                        'img_w': cached['img_w'], 'img_h': cached['img_h'],
                        'datetime_str': cached['datetime_str'],
                        'lat': cached['lat'], 'lon': cached['lon'],
                        'elev_ft': cached.get('elev_ft'), 'elev_m': cached.get('elev_m'),
                        'area_phrase': cached.get('area_phrase', ''),
                        'photographer': cached.get('photographer', ''),
                    }
                    topo_path = FOLDER / f'{stem}_topo.jpg'
                    missing_keys = {k for k, _ in missing}
                    # Restore topo first (uncaptioned) so fullinfo can use it as inset.
                    if 'topo' in missing_keys:
                        print(f'  Fetching topo map...')
                        fetch_topo(cmeta['lat'], cmeta['lon'], topo_path)
                    if 'fullinfo' in missing_keys:
                        print(f'  Regenerating _fullinfo...')
                        make_fullinfo(photo, topo_path, FOLDER / f'{stem}_fullinfo.jpg', cmeta)
                    if 'nearestarea' in missing_keys:
                        print(f'  Regenerating _nearestarea...')
                        make_nearestarea(photo, FOLDER / f'{stem}_nearestarea.jpg', cmeta)
                    # Caption the topo LAST so the clean topo was available as inset above.
                    if 'topo' in missing_keys:
                        make_topo_captioned(topo_path, cmeta)
                else:
                    print(f'Unchanged: {fname}')
                continue

            if sig_matches and cached.get('date_flag'):
                print(f'Re-checking previously flagged ({cached["date_flag"]}): {fname}')
            else:
                print(f'\nProcessing: {fname}')

            # EXIF — cheap (local, no network), so re-run even for a previously-flagged
            # photo whose sig hasn't changed: a code change (e.g. a new
            # CAMERA_TIME_OFFSETS entry, or a wider DATE_TRUST_BUFFER_DAYS) may resolve
            # it now without needing --force.
            print('  Extracting EXIF...')
            exif_raw = extract_exif(photo)
            fallback_dt = datetime.fromtimestamp(photo.stat().st_mtime)
            parsed   = parse_exif(exif_raw, fallback_dt=fallback_dt, fname=fname)

            if parsed['date_flag']:
                # No date we can trust yet — skip elevation/geocode/topo/fullinfo/nearestarea
                # entirely. This photo will only show up in DATE_REVIEW_NEEDED.txt until its
                # date is fixed (in EXIF, via a CAMERA_TIME_OFFSETS entry, or a future manual-
                # override mechanism) and the flag clears on a later run.
                print(f'  ** {parsed["date_flag"]}: {parsed["date_flag_raw_exif"] or "(none)"} '
                      f'— skipping all processing, review-list only **')
                state[fname] = {
                    'sig':                list(sig),
                    'date_flag':          parsed['date_flag'],
                    'date_flag_raw_exif': parsed['date_flag_raw_exif'],
                    'datetime_raw':       parsed['datetime_raw'],
                    'photographer':       parsed['photographer'],
                }
                save_state(state)
                continue

            has_gps = parsed['lat'] is not None

            if has_gps:
                # Elevation
                print('  USGS elevation...')
                elev_ft, elev_m = get_elevation(parsed['lat'], parsed['lon'])

                # Area phrase
                print('  Reverse geocode...')
                area_phrase = get_area_phrase(parsed['lat'], parsed['lon'])

                # Topo image
                topo_fname = f'{stem}_topo.jpg'
                topo_path  = FOLDER / topo_fname
                print(f'  Topo map → {topo_fname}')
                fetch_topo(parsed['lat'], parsed['lon'], topo_path)
            else:
                print('  No GPS — skipping elevation, geocode, topo')
                elev_ft = elev_m = None
                area_phrase = ''
                topo_fname = None
                topo_path  = None

            # Captioned images
            meta = {**parsed, 'elev_ft': elev_ft, 'elev_m': elev_m,
                    'area_phrase': area_phrase}

            fi_fname = f'{stem}_fullinfo.jpg'
            na_fname = f'{stem}_nearestarea.jpg' if has_gps else None
            print(f'  Captions → {fi_fname}' + (f', {na_fname}' if na_fname else '') )
            make_fullinfo(photo, topo_path, FOLDER / fi_fname, meta)
            if has_gps:
                make_nearestarea(photo, FOLDER / na_fname, meta)
                # Caption the topo LAST so the inset in fullinfo used the clean version.
                make_topo_captioned(topo_path, meta)

            # Update state
            state[fname] = {
                'sig':          list(sig),
                'datetime_raw': parsed['datetime_raw'],
                'datetime_str': parsed['datetime_str'],
                'lat':          parsed['lat'],
                'lon':          parsed['lon'],
                'img_w':        parsed['img_w'],
                'img_h':        parsed['img_h'],
                'elev_ft':      elev_ft,
                'elev_m':       elev_m,
                'area_phrase':  area_phrase,
                'photographer': parsed['photographer'],
                'topo':         topo_fname,
                'fullinfo':     fi_fname,
                'nearestarea':  na_fname,
                'date_flag':          parsed['date_flag'],
                'date_flag_raw_exif': parsed['date_flag_raw_exif'],
            }
            save_state(state)
            print(f'  Done: {fname}')

    n_flagged = write_date_review_report(state)
    if n_flagged:
        print(f'\n{n_flagged} photo(s) need a manual date/time review -> {DATE_REVIEW_REPORT.name}')

    build_manifest(state)

    if args.powerpoint:
        build_powerpoint(state)

    # ── Slideshow prompt ──────────────────────────────────────────────────────
    if not args.no_slideshow:
        if args.slideshow:
            build_slideshow(state)
        else:
            print()
            try:
                ans = input('Build 4K slideshow now? (CRF=14, may take several minutes) [y/N]: ').strip().lower()
            except (EOFError, KeyboardInterrupt):
                ans = 'n'
            if ans == 'y':
                build_slideshow(state)
            else:
                print('Slideshow skipped.  Run with --slideshow to encode later.')

    print('\nAll done.')

if __name__ == '__main__':
    main()
